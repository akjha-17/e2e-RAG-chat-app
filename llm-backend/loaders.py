from pathlib import Path
from typing import Tuple, List
from markdown import markdown
from bs4 import BeautifulSoup
import docx
from pptx import Presentation
from openpyxl import load_workbook
import fitz
from langchain_core.documents import Document
import os
from config import DATA_DIR

def _md_to_text_and_images(md: str, base_path: Path) -> List[Document]:
    html = markdown(md, extensions=['extra', 'tables'])
    soup = BeautifulSoup(html, 'html.parser')
    text = soup.get_text(separator=' ')
    images = []
    for img in soup.find_all('img'):
        src = img.get('src')
        if src:
            images.append(str((base_path.parent / src).resolve()))
    return [Document(page_content=text,metadata={"source_file": os.path.basename(base_path), "file_type": "md"})]

def load_text_from_file(path: str) -> List[Document]:
    """
    Returns (plain_text, image_paths). Non-fatal parsing (best effort).
    """
    documents=[]
    suffix = Path(path).suffix.lower()
    if suffix == ".md":
        txt = Path(path).read_text(encoding="utf-8", errors="ignore")
        return _md_to_text_and_images(txt, path)
    if suffix == ".txt":
        return [Document(page_content=Path(path).read_text(encoding="utf-8", errors="ignore"),
                        metadata={"source_file": os.path.basename(path), "file_type": "txt"})]
    if suffix == ".pdf":
        try:
            documents = _extract_text_from_pdf(str(path))
            return documents
        except Exception:
            return documents
    if suffix in [".docx", ".doc"]:
        try:
            d = docx.Document(str(path))
            text = "\n".join(p.text for p in d.paragraphs)
            return [Document(page_content=text,
                            metadata={"source_file": os.path.basename(path), "file_type": "docx"})]
        except Exception:
            return documents
    if suffix == ".pptx":
        try:
            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    if hasattr(shp, "text"):
                        parts.append(shp.text)
            text = "\n".join(parts)
            return [Document(page_content=text,
                            metadata={"source_file": os.path.basename(path), "file_type": "pptx"})]
        except Exception:
            return documents
    if suffix == ".xlsx":
        try:
            wb = load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for name in wb.sheetnames[:3]:
                ws = wb[name]
                for row in ws.iter_rows(min_row=1, max_row=min(200, ws.max_row), values_only=True):
                    parts.append(" ".join("" if c is None else str(c) for c in row))
            text = "\n".join(parts)
            return [Document(page_content=text,
                            metadata={"source_file": os.path.basename(path), "file_type": "xlsx"})]
        except Exception:
            return documents
    # fallback: try text 
    try:
        # optional: skip extremely large files
        file_path = Path(path)
        if file_path.exists() and file_path.stat().st_size > (50 * 1024 * 1024):
            return documents
        return [Document(page_content=file_path.read_text(encoding='utf-8', errors='ignore'),
                         metadata={"source_file": os.path.basename(path), "file_type": "txt"})]
    except Exception:
        return []

def _extract_text_from_pdf(pdf_path: str) -> List[Document]:
    """Extract text from PDF file."""
    doc = fitz.open(pdf_path)
    documents = []
    for i, page in enumerate(doc,start=1):
        text=page.get_text("text")
        if text.strip():  # Only add non-empty pages
            documents.append(
                Document(page_content=text, 
                         metadata={"source_file": os.path.basename(pdf_path), "file_type": "pdf", "page_number": i})
            )
    return documents